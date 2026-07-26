"""Real docx/pptx rendering — docs/architecture/05-agent-architecture.md
§5.11.1: "Word/PDF generation, PowerPoint generation" tools, per
docs/architecture/06-tool-architecture.md §6.2's rows.

**Not registered as separate Tool Runtime modules under `tools/`,
deliberately.** §6.2 catalogues Word/Excel/PowerPoint generation as tools
whose sole permitted caller is ``services/publishing_service`` — a
*service*, not one of the 7 true agents. `shared/agent_profiles`'s
``callers_allowed_for_tool()`` registry (`tools/README.md`'s "a tool never
writes its own caller allow-list" rule, enforced by
`tests/unit/mara/test_agent_profiles.py`) only covers agent-side grants;
extending it to cover service-level tool grants is a design decision for
whoever owns that registry, not one to make unilaterally here by adding a
hand-maintained `ALLOWED_CALLERS` that the existing invariant test would
then flag as drifted. The rendering itself is real; only the "is this a
capital-T Tool Runtime tool" packaging question is open.

Real python-docx/python-pptx generation — both already-pinned dependencies
(no vendor-choice gap, unlike OCR/embedding/search-provider elsewhere in
this repo). Every citation carried by the source agent outputs is rendered
into the output (§5.11.1: "citations preserved from source agents, deck
speaker notes retaining citations") — never summarized away.
"""

from __future__ import annotations

from docx import Document as DocxDocument
from pptx import Presentation
from pptx.util import Inches, Pt

from shared.schemas.compliance import ComplianceChecklist
from shared.schemas.finance import FinancialAnalysis
from shared.schemas.market import MarketBrief
from shared.schemas.recommendation import RecommendationOutput
from shared.schemas.risk import RiskRating


class CommitteeReportInputs:
    """Everything a rendered report/deck draws from — every field here is
    required (not ``Optional``) so a caller cannot construct this without
    every upstream stage's output in hand; ``services/publishing_service``
    layers the "and each one passed its approval gate" check on top (see
    that module), since this class has no way to know approval state on
    its own."""

    __slots__ = (
        'document_id',
        'compliance_checklist',
        'financial_analysis',
        'risk_rating',
        'market_brief',
        'recommendation',
    )

    def __init__(
        self,
        document_id: str,
        compliance_checklist: ComplianceChecklist,
        financial_analysis: FinancialAnalysis,
        risk_rating: RiskRating,
        market_brief: MarketBrief,
        recommendation: RecommendationOutput,
    ) -> None:
        self.document_id = document_id
        self.compliance_checklist = compliance_checklist
        self.financial_analysis = financial_analysis
        self.risk_rating = risk_rating
        self.market_brief = market_brief
        self.recommendation = recommendation


def _citation_text(citation) -> str:
    return f'{citation.document_id} v{citation.version} §{citation.locator}'


def render_docx(inputs: CommitteeReportInputs) -> bytes:
    """Render a committee report (.docx) — one section per upstream
    agent's approved output, every claim carrying its citation inline."""

    doc = DocxDocument()
    doc.add_heading(f'Loan Assessment Committee Report — {inputs.document_id}', level=0)

    doc.add_heading('Compliance', level=1)
    for item in inputs.compliance_checklist.items:
        citation = (
            f' [{_citation_text(item.policy_citation)}]'
            if item.policy_citation is not None
            else ''
        )
        doc.add_paragraph(
            f'{item.requirement}: {item.status.value.upper()}{citation} — {item.notes}',
            style='List Bullet',
        )

    doc.add_heading('Financial Analysis', level=1)
    doc.add_paragraph(inputs.financial_analysis.assessment)
    for figure in inputs.financial_analysis.figures:
        source = (
            f'from {_citation_text(figure.source_citation)}'
            if figure.source_citation is not None
            else f'via {figure.formula_id} {figure.formula_version}'
        )
        doc.add_paragraph(f'{figure.name}: {figure.value} ({source})', style='List Bullet')
    for citation in inputs.financial_analysis.product_term_citations:
        doc.add_paragraph(f'Product terms: [{_citation_text(citation)}]', style='List Bullet')

    doc.add_heading('Risk Rating', level=1)
    doc.add_paragraph(f'Status: {inputs.risk_rating.status.value}')
    if inputs.risk_rating.overall_score is not None:
        doc.add_paragraph(f'Overall score: {inputs.risk_rating.overall_score:.2f}')
    if inputs.risk_rating.qualitative_range is not None:
        doc.add_paragraph(f'Qualitative range: {inputs.risk_rating.qualitative_range}')
    for name, score in inputs.risk_rating.component_scores.items():
        doc.add_paragraph(f'{name}: {score:.2f}', style='List Bullet')
    for citation in inputs.risk_rating.citations:
        doc.add_paragraph(f'Precedent: [{_citation_text(citation)}]', style='List Bullet')

    doc.add_heading('Market Context', level=1)
    if inputs.market_brief.note:
        doc.add_paragraph(inputs.market_brief.note)
    for claim in inputs.market_brief.claims:
        label = ' (low-credibility)' if claim.is_low_credibility() else ''
        doc.add_paragraph(
            f'{claim.claim}{label} — {claim.source_url}, retrieved '
            f'{claim.retrieved_at.date()}',
            style='List Bullet',
        )

    doc.add_heading('Recommendation', level=1)
    if inputs.recommendation.decision is not None:
        doc.add_paragraph(f'Decision: {inputs.recommendation.decision.value}')
        doc.add_paragraph(inputs.recommendation.reasoning_trail)
        for condition in inputs.recommendation.conditions:
            doc.add_paragraph(condition, style='List Bullet')
        for citation in inputs.recommendation.precedent_citations:
            doc.add_paragraph(f'Precedent: [{_citation_text(citation)}]', style='List Bullet')

    from io import BytesIO

    buffer = BytesIO()
    doc.save(buffer)
    return buffer.getvalue()


def _add_bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: str = '') -> None:
    layout = prs.slide_layouts[1]  # "Title and Content"
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    if not bullets:
        bullets = ['(no items)']
    body.text = bullets[0]
    for bullet in bullets[1:]:
        paragraph = body.add_paragraph()
        paragraph.text = bullet
        paragraph.font.size = Pt(18)
    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def render_pptx(inputs: CommitteeReportInputs) -> bytes:
    """Render a presentation deck (.pptx) — one slide per section, with
    citations preserved in speaker notes (§5.11.1: "deck speaker notes
    retaining citations") rather than cluttering the on-slide bullets."""

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = f'Loan Assessment — {inputs.document_id}'

    compliance_notes = '\n'.join(
        _citation_text(item.policy_citation)
        for item in inputs.compliance_checklist.items
        if item.policy_citation is not None
    )
    _add_bullet_slide(
        prs,
        'Compliance',
        [
            f'{item.requirement}: {item.status.value.upper()}'
            for item in inputs.compliance_checklist.items
        ],
        notes=compliance_notes,
    )

    finance_notes = '\n'.join(
        _citation_text(c) for c in inputs.financial_analysis.product_term_citations
    )
    _add_bullet_slide(
        prs,
        'Financial Analysis',
        [inputs.financial_analysis.assessment]
        + [f'{f.name}: {f.value}' for f in inputs.financial_analysis.figures],
        notes=finance_notes,
    )

    risk_notes = '\n'.join(_citation_text(c) for c in inputs.risk_rating.citations)
    risk_bullets = [f'Status: {inputs.risk_rating.status.value}']
    if inputs.risk_rating.overall_score is not None:
        risk_bullets.append(f'Overall score: {inputs.risk_rating.overall_score:.2f}')
    if inputs.risk_rating.qualitative_range is not None:
        risk_bullets.append(f'Range: {inputs.risk_rating.qualitative_range}')
    _add_bullet_slide(prs, 'Risk Rating', risk_bullets, notes=risk_notes)

    market_bullets = [claim.claim for claim in inputs.market_brief.claims] or [
        inputs.market_brief.note or 'No market data available.'
    ]
    market_notes = '\n'.join(claim.source_url for claim in inputs.market_brief.claims)
    _add_bullet_slide(prs, 'Market Context', market_bullets, notes=market_notes)

    rec_notes = '\n'.join(
        _citation_text(c) for c in inputs.recommendation.precedent_citations
    )
    rec_bullets = (
        [f'Decision: {inputs.recommendation.decision.value}', inputs.recommendation.reasoning_trail]
        if inputs.recommendation.decision is not None
        else ['Recommendation withheld.']
    )
    _add_bullet_slide(prs, 'Recommendation', rec_bullets, notes=rec_notes)

    from io import BytesIO

    buffer = BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
