"""Unit tests for services/publishing_service."""

from __future__ import annotations

from datetime import UTC, datetime

import pytest
from docx import Document as DocxDocument
from pptx import Presentation

from services.publishing_service import (
    CommitteeReportInputs,
    MissingApprovedInputError,
    publish_committee_materials,
)
from shared.schemas.compliance import (
    ComplianceChecklist,
    ComplianceChecklistItem,
    ComplianceStatus,
)
from shared.schemas.finance import FigureProvenance, FinancialAnalysis, FinancialFigure
from shared.schemas.knowledge import PolicyCitation
from shared.schemas.market import MarketBrief, MarketClaim
from shared.schemas.recommendation import RecommendationDecision, RecommendationOutput
from shared.schemas.risk import RiskRating, RiskStatus


def _inputs(
    *,
    recommendation_decision: RecommendationDecision | None = RecommendationDecision.APPROVE,
    risk_status: RiskStatus = RiskStatus.COMPLETE,
) -> CommitteeReportInputs:
    compliance = ComplianceChecklist(
        document_id='doc-1',
        items=[
            ComplianceChecklistItem(
                requirement='valid business registration',
                status=ComplianceStatus.PASS,
                policy_citation=PolicyCitation(document_id='pol-1', version='v3', locator='4.2'),
                notes='Confirmed current.',
            )
        ],
    )
    finance = FinancialAnalysis(
        document_id='doc-1',
        figures=[
            FinancialFigure(
                name='dscr',
                value=1.2,
                provenance=FigureProvenance.CALCULATED,
                formula_id='dscr',
                formula_version='v1',
            )
        ],
        assessment='Adequate repayment capacity.',
        assessment_confidence=0.9,
    )
    risk = RiskRating(
        document_id='doc-1',
        status=risk_status,
        component_scores={'financial': 0.2, 'compliance': 0.1, 'market': 0.2}
        if risk_status == RiskStatus.COMPLETE
        else {},
        overall_score=0.18 if risk_status == RiskStatus.COMPLETE else None,
        confidence=0.9,
        missing_inputs=[] if risk_status == RiskStatus.COMPLETE else ['market'],
    )
    market = MarketBrief(
        sector='F&B',
        region='Klang Valley',
        claims=[
            MarketClaim(
                claim='F&B sector growing.',
                source_url='https://trusted-news.my/x',
                retrieved_at=datetime(2026, 1, 1, tzinfo=UTC),
                credibility_score=0.9,
            )
        ],
    )
    recommendation = RecommendationOutput(
        document_id='doc-1',
        decision=recommendation_decision,
        reasoning_trail='Solid case overall.' if recommendation_decision else '',
        confidence=0.9 if recommendation_decision else 0.0,
        withheld_reason=None if recommendation_decision else 'testing withheld path',
    )
    return CommitteeReportInputs(
        document_id='doc-1',
        compliance_checklist=compliance,
        financial_analysis=finance,
        risk_rating=risk,
        market_brief=market,
        recommendation=recommendation,
    )


class TestValidation:
    @pytest.mark.asyncio
    async def test_withheld_recommendation_refuses_to_publish(self) -> None:
        with pytest.raises(MissingApprovedInputError):
            await publish_committee_materials(_inputs(recommendation_decision=None))

    @pytest.mark.asyncio
    async def test_incomplete_risk_refuses_to_publish(self) -> None:
        with pytest.raises(MissingApprovedInputError):
            await publish_committee_materials(_inputs(risk_status=RiskStatus.INCOMPLETE))


class TestRendering:
    @pytest.mark.asyncio
    async def test_renders_real_docx_and_pptx(self) -> None:
        artifacts = await publish_committee_materials(_inputs())

        assert artifacts.docx_bytes[:2] == b'PK'  # docx/pptx are zip containers
        assert artifacts.pptx_bytes[:2] == b'PK'

        from io import BytesIO

        doc = DocxDocument(BytesIO(artifacts.docx_bytes))
        full_text = '\n'.join(p.text for p in doc.paragraphs)
        assert 'pol-1' in full_text  # compliance citation preserved
        assert 'Adequate repayment capacity' in full_text

        prs = Presentation(BytesIO(artifacts.pptx_bytes))
        assert len(list(prs.slides)) >= 5

    @pytest.mark.asyncio
    async def test_deck_generation_can_be_skipped(self) -> None:
        artifacts = await publish_committee_materials(_inputs(), generate_deck=False)

        assert artifacts.docx_bytes is not None
        assert artifacts.pptx_bytes is None

    @pytest.mark.asyncio
    async def test_citations_preserved_in_pptx_speaker_notes(self) -> None:
        artifacts = await publish_committee_materials(_inputs())

        from io import BytesIO

        prs = Presentation(BytesIO(artifacts.pptx_bytes))
        all_notes = '\n'.join(
            slide.notes_slide.notes_text_frame.text
            for slide in prs.slides
            if slide.has_notes_slide
        )
        assert 'pol-1' in all_notes


class TestAuditLogging:
    @pytest.mark.asyncio
    async def test_success_is_logged(self) -> None:
        entries = []
        await publish_committee_materials(_inputs(), audit_sink=entries.append)

        assert len(entries) == 1
        assert entries[0].outcome == 'success'
